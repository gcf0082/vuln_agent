"""Analyze PPTX slides for layout issues programmatically."""
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys

prs = Presentation("/root/projects/vuln_agent/vuln-agent-intro.pptx")

slide_w = prs.slide_width
slide_h = prs.slide_height
SW = slide_w  # in EMU
SH = slide_h
SW_INCH = slide_w / 914400
SH_INCH = slide_h / 914400

print(f"=== Slide dimensions: {SW_INCH:.2f}\" x {SH_INCH:.2f}\" ===")
print()

MIN_MARGIN = Emu(int(0.5 * 914400))  # 0.5 inch
MIN_GAP = Emu(int(0.3 * 914400))     # 0.3 inch

def emu_to_inch(emu):
    return emu / 914400

def get_shape_bounds(shape):
    """Return (left, top, right, bottom) in EMU."""
    l = shape.left
    t = shape.top
    r = l + shape.width
    b = t + shape.height
    return l, t, r, b

def rects_overlap(a, b):
    """Check if two rects overlap (share any area)."""
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    if ax1 >= bx2 or bx1 >= ax2:
        return False
    if ay1 >= by2 or by1 >= ay2:
        return False
    return True

def rect_distance(a, b):
    """Minimum distance between two rects (0 if overlapping)."""
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    # horizontal distance
    h_dist = max(bx1 - ax2, ax1 - bx2, 0)
    v_dist = max(by1 - ay2, ay1 - by2, 0)
    if h_dist == 0 and v_dist == 0:
        # check if overlapping
        if rects_overlap(a, b):
            return 0
        # adjacent
        return 0
    # Euclidean distance between rects
    return (h_dist**2 + v_dist**2)**0.5

def get_text(shape):
    """Get all text from a shape."""
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    texts.append(t)
    return texts

def luminance(r, g, b):
    """Relative luminance."""
    def ch(c):
        c = c / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(b)

def contrast_ratio(c1, c2):
    """WCAG contrast ratio."""
    l1 = luminance(*c1)
    l2 = luminance(*c2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

def shape_color_summary(shape):
    """Try to get fill and font colors from shape."""
    fill_color = None
    font_colors = set()
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    font_colors.add(str(run.font.color.rgb))
    # Check table cells
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            font_colors.add(str(run.font.color.rgb))
    return fill_color, font_colors

def check_low_contrast(shape, bg_color=(255, 255, 255)):
    """Check text for low contrast against expected background."""
    issues = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    fg = (rgb[0], rgb[1], rgb[2])
                    cr = contrast_ratio(fg, bg_color)
                    if cr < 4.5:
                        texts = get_text(shape)
                        sample = texts[0] if texts else ""
                        issues.append(f"  LOW CONTRAST: text '{sample[:40]}' ratio={cr:.2f} fg=#{rgb} on bg=#{bg_color[0]:02x}{bg_color[1]:02x}{bg_color[2]:02x}")
    if shape.has_table:
        for row_idx, row in enumerate(shape.table.rows):
            for col_idx, cell in enumerate(row.cells):
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            rgb = run.font.color.rgb
                            fg = (rgb[0], rgb[1], rgb[2])
                            cr = contrast_ratio(fg, bg_color)
                            if cr < 4.5:
                                issues.append(f"  LOW CONTRAST in table[{row_idx},{col_idx}]: '{cell.text[:30]}' ratio={cr:.2f}")
    return issues

for slide_idx, slide in enumerate(prs.slides):
    print(f"{'='*70}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*70}")
    
    shapes = list(slide.shapes)
    bounds = [get_shape_bounds(s) for s in shapes]
    
    all_issues = []
    
    # 1. Check margins from slide edges
    for i, shape in enumerate(shapes):
        name = shape.name or f"shape_{i}"
        l, t, r, b = bounds[i]
        texts = get_text(shape)
        sample = texts[0] if texts else ""
        
        if l < MIN_MARGIN:
            all_issues.append(f"  MARGIN: Left edge of '{name}' ({sample[:30]}) at {emu_to_inch(l):.2f}\" < 0.5\"")
        if t < MIN_MARGIN:
            all_issues.append(f"  MARGIN: Top edge of '{name}' ({sample[:30]}) at {emu_to_inch(t):.2f}\" < 0.5\"")
        if SW - r < MIN_MARGIN:
            all_issues.append(f"  MARGIN: Right edge of '{name}' ({sample[:30]}) at {emu_to_inch(SW - r):.2f}\" from right < 0.5\"")
        if SH - b < MIN_MARGIN:
            all_issues.append(f"  MARGIN: Bottom edge of '{name}' ({sample[:30]}) at {emu_to_inch(SH - b):.2f}\" from bottom < 0.5\"")
    
    # 2. Check overlaps
    for i in range(len(shapes)):
        for j in range(i+1, len(shapes)):
            if rects_overlap(bounds[i], bounds[j]):
                ti = get_text(shapes[i])
                tj = get_text(shapes[j])
                si = ti[0] if ti else shapes[i].name
                sj = tj[0] if tj else shapes[j].name
                all_issues.append(f"  OVERLAP: '{si[:30]}' and '{sj[:30]}'")
    
    # 3. Check proximity (too close but not overlapping)
    for i in range(len(shapes)):
        for j in range(i+1, len(shapes)):
            # only check if both have text or are visible
            a_left, a_top, a_right, a_bottom = bounds[i]
            b_left, b_top, b_right, b_bottom = bounds[j]
            
            # vertical gap
            v_gap = max(b_top - a_bottom, a_top - b_bottom)
            # horizontal gap  
            h_gap = max(b_left - a_right, a_left - b_right)
            
            # If they're in same vertical column (overlap in x)
            x_overlap = min(a_right, b_right) - max(a_left, b_left)
            y_overlap = min(a_bottom, b_bottom) - max(a_top, b_top)
            
            if x_overlap > 0 and y_overlap > 0:
                continue  # already reported as overlap
                
            if x_overlap > 0 and 0 < v_gap < MIN_GAP:
                ti = get_text(shapes[i])
                tj = get_text(shapes[j])
                si = ti[0][:30] if ti else shapes[i].name
                sj = tj[0][:30] if tj else shapes[j].name
                all_issues.append(f"  TOO CLOSE (vert): '{si}' and '{sj}' gap={emu_to_inch(v_gap):.2f}\"")
            
            if y_overlap > 0 and 0 < h_gap < MIN_GAP:
                ti = get_text(shapes[i])
                tj = get_text(shapes[j])
                si = ti[0][:30] if ti else shapes[i].name
                sj = tj[0][:30] if tj else shapes[j].name
                all_issues.append(f"  TOO CLOSE (horiz): '{si}' and '{sj}' gap={emu_to_inch(h_gap):.2f}\"")
    
    # 4. Check text overflow in shapes
    for i, shape in enumerate(shapes):
        if shape.has_text_frame:
            tf = shape.text_frame
            # Check if text box is very narrow
            width_in = emu_to_inch(shape.width)
            if width_in < 1.5 and tf.text.strip():
                longest_word = max(tf.text.split(), key=len) if tf.text.split() else ""
                # rough check: if longest word > width, could overflow
                all_issues.append(f"  NARROW TEXTBOX: '{shape.name}' width={width_in:.2f}\" text='{tf.text[:40]}'")
            
            # Check if text might overflow vertically
            # Count lines
            total_lines = 0
            for para in tf.paragraphs:
                total_lines += 1
                # rough estimate: if text is long and width is small
                if para.text and width_in > 0:
                    approx_chars_per_line = int(width_in * 12)  # ~12 chars per inch at ~14pt
                    if approx_chars_per_line > 0:
                        extra_lines = len(para.text) // approx_chars_per_line
                        total_lines += max(0, extra_lines - 1)
            
            font_size = None
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size = run.font.size
                        break
                if font_size:
                    break
            
            if font_size:
                line_height = font_size * 1.3  # approximate
                needed_height = total_lines * line_height
                available_height = shape.height
                if needed_height > available_height * 1.1:  # 10% tolerance
                    all_issues.append(f"  TEXT OVERFLOW: '{shape.name}' needs ~{emu_to_inch(needed_height):.2f}\" but box is {emu_to_inch(available_height):.2f}\"")
    
    # 5. Check for placeholder/decorative lines
    for i, shape in enumerate(shapes):
        texts = get_text(shape)
        if not texts and not shape.has_table:
            # Could be a line or shape
            if shape.width < Emu(50000):  # thin = likely a line
                all_issues.append(f"  DECORATIVE LINE: '{shape.name}' at y={emu_to_inch(shape.top):.2f}\"")
    
    # 6. Check alignment consistency - group shapes by approximate x position
    # Group shapes by left position (within 0.2")
    col_groups = {}
    for i, shape in enumerate(shapes):
        left_in = round(emu_to_inch(shape.left), 1)
        if left_in not in col_groups:
            col_groups[left_in] = []
        col_groups[left_in].append(i)
    
    if len(col_groups) > 1:
        for col_x, members in col_groups.items():
            if len(members) > 1:
                lefts = [shapes[m].left for m in members]
                widths = [shapes[m].width for m in members]
                # Check if same-column shapes have consistent widths
                if len(set(widths)) > 1 and max(widths) - min(widths) > Emu(int(0.2 * 914400)):
                    texts_in_col = [get_text(shapes[m])[0][:30] if get_text(shapes[m]) else shapes[m].name for m in members]
                    all_issues.append(f"  INCONSISTENT WIDTH in column x={col_x}\": {texts_in_col}")
    
    # 7. Check table-specific issues
    for i, shape in enumerate(shapes):
        if shape.has_table:
            tbl = shape.table
            # Check column widths consistency
            col_widths = [col.width for col in tbl.columns]
            if col_widths and len(col_widths) > 1:
                for w in col_widths:
                    if w < Emu(int(0.5 * 914400)):
                        all_issues.append(f"  NARROW TABLE COLUMN: width={emu_to_inch(w):.2f}\"")
            
            # Check for cramped cells
            for row_idx, row in enumerate(tbl.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    if cell_text:
                        # Check if cell is very small
                        cell_w = tbl.columns[col_idx].width
                        if cell_w < Emu(int(0.6 * 914400)) and len(cell_text) > 10:
                            all_issues.append(f"  TIGHT TABLE CELL [{row_idx},{col_idx}]: '{cell_text[:30]}' in {emu_to_inch(cell_w):.2f}\" wide cell")
    
    # 8. Check for placeholder text
    for i, shape in enumerate(shapes):
        texts = get_text(shape)
        for t in texts:
            tl = t.lower()
            if any(x in tl for x in ['xxxx', 'lorem', 'ipsum', 'placeholder', 'click to add']):
                all_issues.append(f"  PLACEHOLDER TEXT: '{t}' in '{shape.name}'")
    
    # 9. Low contrast check (assume white/cream background)
    for i, shape in enumerate(shapes):
        issues = check_low_contrast(shape, bg_color=(255, 255, 255))
        all_issues.extend(issues)
        # Also check against cream
        issues2 = check_low_contrast(shape, bg_color=(245, 240, 230))
        all_issues.extend(issues2)
    
    # Remove duplicates
    all_issues = list(dict.fromkeys(all_issues))
    
    if all_issues:
        for issue in all_issues:
            print(issue)
    else:
        print("  No issues detected.")
    
    # 10. Print shape inventory for manual review
    print()
    print("  Shape inventory:")
    for i, shape in enumerate(shapes):
        l = emu_to_inch(shape.left)
        t = emu_to_inch(shape.top)
        w = emu_to_inch(shape.width)
        h = emu_to_inch(shape.height)
        r = emu_to_inch(shape.left + shape.width)
        b = emu_to_inch(shape.top + shape.height)
        texts = get_text(shape)
        sample = texts[0][:50] if texts else "(no text)"
        shape_type = shape.shape_type
        print(f"  [{i}] {shape.name} type={shape_type} pos=({l:.2f},{t:.2f}) size=({w:.2f}x{h:.2f}) right={r:.2f} bot={b:.2f} text='{sample}'")
    print()
